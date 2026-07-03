#!/usr/bin/env python3
"""
Generate three InsideHR training decks (admin, manager, employee) as .pptx
files in docs/decks/.

Style: iORA-inspired — black on white, minimalist.
Screenshots are read from docs/screenshots/.

Run:  python3 scripts/generate_decks.py
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
SCREENS = ROOT / "docs" / "screenshots"
OUT_DIR = ROOT / "docs" / "decks"

# ---- Colours (iORA palette) ----
INK = RGBColor(0x18, 0x18, 0x1B)        # primary text / titles
INK_SOFT = RGBColor(0x3F, 0x3F, 0x46)   # body text
MUTE = RGBColor(0x71, 0x71, 0x7A)       # captions
RULE = RGBColor(0xE4, 0xE4, 0xE7)       # divider lines
ACCENT_NAVY = RGBColor(0x1E, 0x3A, 0x5F)
BLUSH = RGBColor(0xE8, 0xB4, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_TINT = RGBColor(0xFA, 0xFA, 0xFA)

# ---- Layout (16:9, 13.333" x 7.5") ----
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.55)


@dataclass
class Slide:
    """Slide spec — title + optional subtitle + optional screenshot + steps/bullets."""

    title: str
    subtitle: Optional[str] = None
    screenshot: Optional[str] = None  # filename in docs/screenshots/, without .png
    steps: Optional[List[str]] = None
    bullets: Optional[List[str]] = None
    note: Optional[str] = None
    section: bool = False  # render as a divider/section slide


# ============================================================
# Helpers
# ============================================================


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 14,
    color: RGBColor = INK_SOFT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    line_spacing: float = 1.25,
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold


def add_rect(slide, left, top, width, height, fill: RGBColor, line: Optional[RGBColor] = None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_line(slide, x1, y1, x2, y2, color=RULE, weight=0.5):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def image_max_size(path: Path, max_w: Emu, max_h: Emu) -> Tuple[Emu, Emu]:
    """Fit-inside dimensions preserving aspect ratio."""
    with Image.open(path) as im:
        w, h = im.size
    img_aspect = w / h
    box_aspect = max_w / max_h
    if img_aspect > box_aspect:
        target_w = max_w
        target_h = int(max_w / img_aspect)
    else:
        target_h = max_h
        target_w = int(max_h * img_aspect)
    return Emu(target_w), Emu(target_h)


def add_screenshot(slide, name: str, *, left, top, max_w, max_h, caption: Optional[str] = None):
    path = SCREENS / f"{name}.png"
    if not path.exists():
        # Placeholder if the asset is missing
        add_rect(slide, left, top, max_w, max_h, BG_TINT, line=RULE)
        add_text(
            slide, left, top + Inches(0.4),
            max_w, Inches(0.4),
            f"[Screenshot missing: {name}.png]",
            size=12, color=MUTE, align=PP_ALIGN.CENTER,
        )
        return

    w, h = image_max_size(path, max_w, max_h)
    # Centre inside the bounding box
    cx = left + int((max_w - w) / 2)
    cy = top + int((max_h - h) / 2)
    pic = slide.shapes.add_picture(str(path), cx, cy, width=w, height=h)
    # Thin border
    pic.line.color.rgb = RULE
    pic.line.width = Pt(0.5)
    if caption:
        add_text(
            slide, left, top + max_h + Inches(0.05),
            max_w, Inches(0.3),
            caption, size=10, color=MUTE, align=PP_ALIGN.CENTER,
        )


# ============================================================
# Slide builders
# ============================================================


def build_title_slide(prs, *, role: str, sub: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Wordmark top-left
    add_text(slide, MARGIN, Inches(0.5), Inches(3), Inches(0.5),
             "InsideHR", size=18, color=INK, bold=True)
    # Big title
    add_text(slide, MARGIN, Inches(2.4), Inches(12), Inches(1.4),
             f"{role} Training", size=54, color=INK, bold=True)
    # Subtitle
    add_text(slide, MARGIN, Inches(3.8), Inches(12), Inches(0.6),
             sub, size=20, color=MUTE)
    # Hairline rule
    add_line(slide, MARGIN, Inches(4.7), Inches(5.5), Inches(4.7), color=INK, weight=1.5)
    # Footer
    add_text(slide, MARGIN, Inches(6.9), Inches(12), Inches(0.3),
             "An iORA-style HR system for Singapore + Malaysia",
             size=11, color=MUTE)


def build_section_slide(prs, *, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Black band on the left
    add_rect(slide, 0, 0, Inches(0.7), SLIDE_H, INK)
    # Section number / title
    add_text(slide, Inches(1.2), Inches(2.9), Inches(11), Inches(1.2),
             title, size=44, color=INK, bold=True)
    add_text(slide, Inches(1.2), Inches(4.2), Inches(11), Inches(0.6),
             subtitle, size=18, color=MUTE)


def build_content_slide(prs, spec: Slide, *, deck_name: str, page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header band: deck label + page number
    add_text(slide, MARGIN, Inches(0.32), Inches(8), Inches(0.3),
             deck_name, size=10, color=MUTE)
    add_text(slide, Inches(11.5), Inches(0.32), Inches(1.5), Inches(0.3),
             f"{page_num}", size=10, color=MUTE, align=PP_ALIGN.RIGHT)
    add_line(slide, MARGIN, Inches(0.78), SLIDE_W - MARGIN, Inches(0.78))

    # Title
    add_text(slide, MARGIN, Inches(0.95), Inches(12), Inches(0.6),
             spec.title, size=28, color=INK, bold=True)

    # Subtitle
    title_h = Inches(0.6)
    sub_y = Inches(0.95) + title_h
    if spec.subtitle:
        add_text(slide, MARGIN, sub_y, Inches(12), Inches(0.5),
                 spec.subtitle, size=14, color=MUTE)
        content_top = sub_y + Inches(0.55)
    else:
        content_top = sub_y + Inches(0.05)

    # Two-column layout if screenshot present: text left, image right
    if spec.screenshot:
        text_left = MARGIN
        text_w = Inches(5.6)
        text_top = content_top
        text_h = SLIDE_H - text_top - Inches(0.7)

        img_left = Inches(6.4)
        img_top = content_top
        img_w = SLIDE_W - img_left - MARGIN
        img_h = text_h

        _render_text_column(slide, spec, text_left, text_top, text_w, text_h)
        add_screenshot(slide, spec.screenshot, left=img_left, top=img_top, max_w=img_w, max_h=img_h)
    else:
        # Full-width text
        _render_text_column(slide, spec, MARGIN, content_top, SLIDE_W - 2 * MARGIN, SLIDE_H - content_top - Inches(0.7))

    # Footer note
    if spec.note:
        add_text(slide, MARGIN, Inches(7.05), SLIDE_W - 2 * MARGIN, Inches(0.35),
                 spec.note, size=10, color=MUTE)


def _render_text_column(slide, spec: Slide, left, top, width, height):
    y = top
    if spec.steps:
        add_text(slide, left, y, width, Inches(0.3),
                 "STEP-BY-STEP", size=10, color=MUTE, bold=True)
        y += Inches(0.32)
        for i, s in enumerate(spec.steps, start=1):
            # Number badge
            badge = add_rect(slide, left, y + Inches(0.05), Inches(0.32), Inches(0.32), INK)
            badge.line.fill.background()
            tb = slide.shapes.add_textbox(left, y + Inches(0.05), Inches(0.32), Inches(0.32))
            tf = tb.text_frame
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(i)
            r.font.name = "Arial"
            r.font.size = Pt(13)
            r.font.color.rgb = WHITE
            r.font.bold = True

            # Step text — wrap-aware textbox to the right of the badge
            step_text_left = left + Inches(0.45)
            step_text_w = width - Inches(0.45)
            tb = slide.shapes.add_textbox(step_text_left, y, step_text_w, Inches(0.8))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.line_spacing = 1.2
            r = p.add_run()
            r.text = s
            r.font.name = "Arial"
            r.font.size = Pt(13)
            r.font.color.rgb = INK_SOFT
            # Rough vertical advance — assume ~ 1 line per 35 chars
            lines = max(1, (len(s) // 60) + 1)
            y += Inches(0.34) + Inches(0.22 * (lines - 1)) + Inches(0.2)

    if spec.bullets:
        add_text(slide, left, y, width, Inches(0.3),
                 "WHAT YOU CAN DO", size=10, color=MUTE, bold=True)
        y += Inches(0.32)
        for b in spec.bullets:
            # Dot
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.05), y + Inches(0.13), Inches(0.08), Inches(0.08))
            dot.fill.solid()
            dot.fill.fore_color.rgb = INK
            dot.line.fill.background()

            tb = slide.shapes.add_textbox(left + Inches(0.25), y, width - Inches(0.25), Inches(0.5))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.line_spacing = 1.2
            r = p.add_run()
            r.text = b
            r.font.name = "Arial"
            r.font.size = Pt(13)
            r.font.color.rgb = INK_SOFT
            lines = max(1, (len(b) // 60) + 1)
            y += Inches(0.34) + Inches(0.22 * (lines - 1))


# ============================================================
# Deck content
# ============================================================


# ---- Admin deck ----
ADMIN_SLIDES: List[Slide] = [
    Slide(
        title="Welcome — Admin role",
        subtitle="The full picture of what InsideHR lets you do as the system owner.",
        bullets=[
            "InsideHR is the people + payroll system for SG and MY operations.",
            "As Admin, you can manage everyone (people, holidays, leave policy), see every department, and configure the new modules (Performance, Timesheets, Rewards).",
            "Managers see only their direct reports. Employees see only their own records.",
            "This deck walks every page step-by-step. ~35 slides; should take ~30 minutes.",
        ],
        note="Roles: ADMIN (full access) · MANAGER (their team) · EMPLOYEE (self-service) · CONTRACTOR · PART_TIME",
        section=True,
    ),
    Slide(
        title="Logging in",
        subtitle="Every staff member uses the same login screen.",
        steps=[
            "Open the InsideHR URL in your browser.",
            "Enter the email + temporary password issued to you.",
            "On first login you'll be forced to change your password (8+ characters).",
            "After that, the password lives only in your head — Admin can reset but never see it.",
        ],
        bullets=[
            "Forgot password? Click the link, and you'll get a reset email via Resend.",
            "The system enforces secure session cookies that expire after a period of inactivity.",
        ],
        note="On first sign-in for the demo: jin@company.com / test123",
    ),
    # SECTION: Dashboard + people
    Slide(title="Dashboard & People", subtitle="Day-to-day overview and the employee directory.", section=True),
    Slide(
        title="The admin dashboard",
        subtitle="At-a-glance summary of what needs your attention.",
        screenshot="admin-dashboard",
        bullets=[
            "Top: approval-count widget — pending leave + expense + reviews across the company.",
            "Birthday widget for the next 14 days.",
            "Country holiday card — defaults to your country's holidays for the current year.",
            "From here you can navigate to any module via the sidebar.",
        ],
    ),
    Slide(
        title="People — directory",
        subtitle="One row per active employee, with quick filter + search.",
        screenshot="admin-people-list",
        bullets=[
            "Filter by country (SG / MY), role, department, status.",
            "Coloured avatar + role pill makes status visible.",
            "Click any row → detailed employee profile with documents, leave, reviews, audit trail.",
        ],
    ),
    Slide(
        title="People — add new employee",
        subtitle="The form that creates a user record.",
        screenshot="admin-people-new",
        steps=[
            "Click \"+ New employee\" on the People list.",
            "Fill out basics (name, email, country), then the job info section.",
            "Set Employment Type — EMPLOYEE / CONTRACTOR / PART_TIME. Part-time enables timesheets + payroll fields.",
            "Pick a reporting manager — drives review assignment, leave approval, and expense approval.",
            "Save → the system issues a temporary password and the user gets an invite email.",
        ],
        note="PART_TIME = sets hourlyRate + normalDailyHours required for payroll math.",
    ),
    Slide(
        title="Org chart",
        subtitle="The reporting hierarchy as a navigable tree.",
        screenshot="admin-org-chart",
        bullets=[
            "Auto-built from each employee's reportingManagerId.",
            "Click a node to focus on that branch.",
            "Useful for sanity-checking the management graph before running a perf review cycle.",
        ],
    ),
    # SECTION: Time off + holidays
    Slide(title="Leave & Holidays", subtitle="Time off policy, balances, public holidays.", section=True),
    Slide(
        title="Holidays — SG & MY",
        subtitle="Public holiday calendar per country.",
        screenshot="admin-holidays",
        bullets=[
            "Seeded with 2026 federal holidays for both countries.",
            "Add / edit / mark observed status per country.",
            "Drives auto-detection in the part-time timesheet (logging hours on a holiday = 2× pay).",
        ],
    ),
    Slide(
        title="Admin Leave Management",
        subtitle="Operations admin: bulk balance adjust, CSV import, year-end carry forward.",
        screenshot="admin-leave-management",
        bullets=[
            "Balance adjust — give an employee +N days for a leave type, audit-logged.",
            "CSV import — bulk-load historical leave from your previous system.",
            "Carry-forward — at year-end, roll unused balance into the new year (caps configurable).",
        ],
    ),
    Slide(
        title="Time Off — view from admin",
        subtitle="Admin can see and act on any employee's leave.",
        screenshot="admin-leave",
        bullets=[
            "Filter by status, leave type, person, or date range.",
            "Admin override: approve / reject any request (audit-logged).",
        ],
    ),
    # SECTION: Documents + expenses
    Slide(title="Documents & Expenses", subtitle="Files and reimbursable spend.", section=True),
    Slide(
        title="Documents",
        subtitle="Company-wide files + per-employee documents.",
        screenshot="admin-documents",
        bullets=[
            "Company tab: HR uploads handbook, policies, holiday calendar PDFs — visible to everyone.",
            "Employee tab: per-employee folders — only visible to that employee + their manager + admin.",
            "Files live in Google Drive; the app issues short-lived download URLs.",
        ],
    ),
    Slide(
        title="Expenses — admin reimbursement queue",
        subtitle="Approved expenses awaiting reimbursement.",
        screenshot="admin-expenses-approvals",
        bullets=[
            "Bulk-select rows, then mark as reimbursed — sets reimbursedAt + reimbursedById, audit-logged.",
            "Export Excel: per-employee summary + line items for finance handoff.",
            "Receipts are linked Google Drive URLs.",
        ],
    ),
    Slide(
        title="Expenses — your own",
        subtitle="Admin can also file expenses like any employee.",
        screenshot="admin-expenses",
        bullets=[
            "Tabs: Draft / For Approval / Approved / Rejected / Reimbursed.",
            "Click \"+ New expense\" → category, amount + currency, merchant, receipt date, description, receipt upload.",
            "Multi-step approval if configured (manager → admin reimburses).",
        ],
    ),
    # SECTION: Performance review — the heart of this deck
    Slide(title="Performance Reviews", subtitle="Cycles, scope, goal-setting, evaluation, lifecycle.", section=True),
    Slide(
        title="Performance — the model",
        subtitle="Three templates, one data model.",
        bullets=[
            "FULL: store managers + HQ. Goals + 5-point rating + narrative.",
            "LITE: part-timers + casuals. Behavioural 3-point rating, no goals. Promotion-ready flag.",
            "PROBATION: 3-month confirmation review for new hires (MY-compliance lever). Decision: CONFIRMED / EXTENDED / NOT_CONFIRMED.",
            "Per-cycle config: rating labels, goal min/max, employee comment, sales target, attendance, etc.",
        ],
        note="Cadence is up to you — create a new cycle quarterly, biannually, or ad-hoc.",
    ),
    Slide(
        title="Step 1 — Create a review cycle",
        subtitle="Set the rules for this round of reviews.",
        screenshot="admin-perf-cycle-new",
        steps=[
            "Sidebar → Performance → Cycles → + New cycle.",
            "Name + start/end dates + (optional) goal-setting & evaluation deadlines.",
            "Pick a template. The form below auto-adjusts (no rating fields for PROBATION; no goals for LITE).",
            "Toggle the retail extras you care about — sales target (with currency), attendance metric.",
            "Save → you land on the cycle detail page (status = DRAFT).",
        ],
    ),
    Slide(
        title="Step 2 — Scope the cycle",
        subtitle="Pick who is in this cycle.",
        screenshot="admin-perf-cycle-detail",
        steps=[
            "On the cycle detail page, find the Scope section.",
            "Filter mode: pick employment type / country / department → click \"Add matching employees\".",
            "Picker mode (below): search the list of active employees not yet in this cycle → tick checkboxes → \"Add picked\".",
            "Each scoped employee gets a PerformanceReview row, with their reporting manager snapshotted at this moment.",
        ],
        note="Manager-snapshot means reorgs mid-cycle don't lose history.",
    ),
    Slide(
        title="Step 3 — Open the cycle",
        subtitle="DRAFT → ACTIVE: managers can now set goals.",
        screenshot="admin-perf-cycle-detail",
        steps=[
            "Lifecycle panel on the cycle detail page.",
            "Click \"Open cycle (start goal-setting)\" — cycle goes ACTIVE.",
            "Each scoped review goes from NOT_STARTED → managers can now write goals for their reports.",
            "Once minGoals are saved per review, that review auto-promotes to GOALS_SET.",
        ],
    ),
    Slide(
        title="Step 4 — Move to evaluation",
        subtitle="ACTIVE → EVALUATION: managers rate + submit.",
        bullets=[
            "Click \"Move to evaluation\" when the goal-setting window closes.",
            "Goals become read-only; each goal needs an outcome (Missed / Partial / Met / Exceeded) plus an optional comment.",
            "Managers also fill the overall rating, narrative, sales target/actual, attendance — depending on what the cycle includes.",
            "Once a manager submits, the review → PENDING_ACKNOWLEDGEMENT.",
        ],
    ),
    Slide(
        title="Step 5 — Watch progress, then close",
        subtitle="EVALUATION → CLOSED.",
        bullets=[
            "Cycle detail shows a table of every review with its status pill.",
            "Once all employees have acknowledged, click \"Close cycle\" → CLOSED, locked from further edits.",
            "Admin can still \"Reopen\" an individual review if a correction is needed (audit-logged).",
        ],
    ),
    Slide(
        title="HR export for closed cycles",
        subtitle="Excel workbook with 5 sheets, ready for HR/payroll.",
        bullets=[
            "Cycle — config snapshot.",
            "Reviews — one row per employee with rating, narrative, sales numbers, probation decision.",
            "Goals — every goal across the cycle with outcomes + manager comments.",
            "Rating distribution — for calibration.",
            "Probation outcomes — when applicable.",
        ],
        note="Click \"Export Excel\" in the cycle header.",
    ),
    # SECTION: Timesheet + payroll
    Slide(title="Part-time Timesheets & Payroll", subtitle="The new module for hourly workers.", section=True),
    Slide(
        title="How the part-time flow works",
        subtitle="Employee submits hours → manager approves → admin runs payroll.",
        bullets=[
            "Part-time employees log hours per day, then submit the whole week.",
            "Submitted entries lock from edit; manager sees them in their approval queue.",
            "Manager approves (or rejects with a reason). Rejected entries unlock for the employee to fix.",
            "Approved entries flow into the monthly payroll calculation.",
        ],
        note="Malaysian Employment Act 1955 multipliers apply automatically (1.5× over 8h/day or 45h/week, 2× PH, 3× PH OT).",
    ),
    Slide(
        title="Time approvals (admin override)",
        subtitle="Admin can also approve directly.",
        screenshot="admin-time-approvals",
        bullets=[
            "Same view as the manager approval queue.",
            "Useful when a manager is on leave or terminated mid-week.",
            "Bulk-approve a whole week per employee.",
        ],
    ),
    Slide(
        title="Payroll — monthly breakdown",
        subtitle="Auto-computed pay per part-timer.",
        screenshot="admin-payroll",
        bullets=[
            "Per-row split: Regular / OT (1.5×) / PH (2×) / PH-OT (3×) hours × hourly rate.",
            "Currency derived from each employee's country (SGD / MYR).",
            "Prev / Next month buttons; defaults to current month.",
            "Export Excel: summary sheet + line items for your payroll provider.",
        ],
        note="Multipliers per Employment Act 1955 (2022 amendment).",
    ),
    # SECTION: Rewards
    Slide(title="Rewards & Bonuses", subtitle="Tie performance to bonus payouts.", section=True),
    Slide(
        title="Rewards — the model",
        subtitle="Discretionary bonus cycles, optionally linked to a perf review.",
        screenshot="admin-rewards-cycles",
        bullets=[
            "Bonus types: PERFORMANCE / CONTRACTUAL_13TH / AD_HOC.",
            "Cycle currency: MYR / SGD / USD.",
            "Lifecycle: DRAFT → APPROVED → PAID → CLOSED — when you transition, all draft allocations cascade with you.",
        ],
    ),
    Slide(
        title="Step 1 — Create a reward cycle",
        screenshot="admin-rewards-cycle-new",
        steps=[
            "Sidebar → Admin → Rewards → + New cycle.",
            "Name (e.g. \"Hari Raya 2026 Bonus\"), optional description.",
            "Optionally link to a perf review cycle — when linked, you'll see each employee's overall rating during allocation, for context.",
            "Pick currency, optional total pool (for budgeting), optional payout date.",
            "Save → land on the cycle detail page in DRAFT.",
        ],
    ),
    Slide(
        title="Step 2 — Allocate amounts",
        subtitle="Per-employee, per-bonus-type allocations.",
        screenshot="admin-rewards-cycle-detail",
        steps=[
            "Expand \"+ Add allocation\" — see all active employees with their linked rating.",
            "Click \"Allocate\" next to one — inline form: bonus type, amount, rationale.",
            "Each row in the table = one allocation with a status pill.",
            "Edit (pencil) or cancel (X) draft allocations any time.",
        ],
    ),
    Slide(
        title="Step 3 — Approve, pay, close",
        subtitle="Move the whole cycle through its lifecycle.",
        bullets=[
            "\"Approve all draft allocations\" → cycle APPROVED, draft rows stamped with you as approver, locked.",
            "\"Mark as paid\" → cycle PAID, all allocations get paidAt timestamp.",
            "\"Close cycle\" → final archive.",
            "Excel export available at any point: Allocations sheet + Summary by bonus type.",
        ],
    ),
    # SECTION: Audit + permissions
    Slide(title="Audit & Permissions", subtitle="Trace and trust.", section=True),
    Slide(
        title="Audit log",
        subtitle="Every state transition leaves a trail.",
        bullets=[
            "Captured per record: who, when, what action, IP.",
            "Includes login, password changes, leave/expense state changes, review submissions, probation outcomes, timesheet submissions, reward approvals.",
            "Visible per-employee on their profile page (recent activity).",
            "Useful for compliance + dispute resolution.",
        ],
    ),
    Slide(
        title="Permissions in one slide",
        subtitle="Who can do what.",
        bullets=[
            "ADMIN — every page, every action. Can reopen acknowledged reviews; can unlock approved timesheet entries; can transition reward cycles.",
            "MANAGER — sees their direct reports. Approves leave / expenses for them; sets and evaluates their goals; approves their timesheet entries.",
            "EMPLOYEE — own records only. Submit leave / expenses; acknowledge perf reviews; (if PART_TIME) log + submit timesheets.",
            "CONTRACTOR — like EMPLOYEE but excluded from default scoping in some flows.",
        ],
    ),
    Slide(
        title="That's the system",
        subtitle="Welcome to InsideHR.",
        bullets=[
            "Questions? Ask your HR Admin or IT.",
            "Bugs? File a ticket with what you were doing + a screenshot.",
            "Feedback on features? The system is evolving — your input shapes it.",
        ],
        note="iORA-style HR for SG + MY operations · v0.1",
    ),
]


# ---- Manager deck ----
MANAGER_SLIDES: List[Slide] = [
    Slide(
        title="Welcome — Manager role",
        subtitle="What you can do for your direct reports.",
        bullets=[
            "As a manager, you see your direct reports' records: leave requests, expense submissions, performance reviews, and timesheets (if you have part-timers).",
            "You approve or reject those requests, you write goals + ratings, you sign off weekly hours.",
            "You can't see other managers' teams or system-wide settings.",
            "This deck walks each thing step-by-step. ~25 slides; ~20 minutes.",
        ],
    ),
    Slide(
        title="Logging in",
        steps=[
            "Open the InsideHR URL.",
            "Enter your work email + the temporary password Admin sent you.",
            "On first sign-in you'll be forced to change your password.",
            "After that, you land on your manager dashboard.",
        ],
        note="If your email isn't recognised, contact Admin — your account hasn't been created yet.",
    ),
    Slide(
        title="Your dashboard",
        subtitle="What needs your attention today.",
        screenshot="manager-dashboard",
        bullets=[
            "Approval-count widget — pending leave + expense requests from your reports.",
            "Birthday widget — direct reports' upcoming birthdays.",
            "Use the sidebar to jump to any specific module.",
        ],
    ),
    # SECTION: Leave
    Slide(title="Leave Approvals", subtitle="Your reports' time-off requests.", section=True),
    Slide(
        title="The leave approval queue",
        screenshot="manager-leave-approvals",
        bullets=[
            "One card per pending leave request from your direct reports.",
            "Each card shows: who, dates, leave type, number of working days, balance left.",
            "Approve immediately, or Reject with an optional comment.",
            "Rejection sends the employee an email with your reason.",
        ],
    ),
    Slide(
        title="Reviewing a leave request",
        steps=[
            "Click into the request from the queue.",
            "See the full reason + any attachment (e.g. medical certificate).",
            "Use the \"Approve\" or \"Reject\" buttons. Rejection requires a reason.",
            "Approved leave deducts from the employee's balance immediately and shows on the team calendar.",
        ],
        note="If you're unsure, ask the employee directly first — the system isn't your manager.",
    ),
    # SECTION: Expenses
    Slide(title="Expense Approvals", subtitle="Reimbursable spend by your reports.", section=True),
    Slide(
        title="Expense approval queue",
        screenshot="manager-expense-approvals",
        bullets=[
            "Tab: For Approval — expenses your reports submitted, pending your sign-off.",
            "Each row: category, amount + currency, receipt date, merchant.",
            "Click into one → see the full receipt (image preview or PDF download).",
            "Approve (passes to admin for reimbursement) or Reject (back to employee to fix).",
        ],
    ),
    # SECTION: Performance review — the meat of this deck
    Slide(title="Performance Reviews", subtitle="Setting goals → evaluating → submitting.", section=True),
    Slide(
        title="Performance — your team queue",
        screenshot="manager-perf-team",
        bullets=[
            "When an Admin scopes you into a cycle, you see your reports here.",
            "Status pill tells you what's needed: \"Not started\" / \"Goals set\" / \"In evaluation\" / \"Awaiting employee\" / \"Acknowledged\".",
            "The action link (\"Set goals →\" / \"Evaluate →\" / \"View →\") changes with cycle state.",
        ],
    ),
    Slide(
        title="Step 1 — Set goals",
        subtitle="During the ACTIVE phase of a cycle.",
        screenshot="manager-perf-review-detail",
        steps=[
            "Click into a report's review → \"+ Add goal\".",
            "Title (1 sentence) + description (the why) + type (qualitative or quantitative).",
            "For quantitative goals: set a target value + unit (RM, units, %, NPS, etc.).",
            "Hit min goals (usually 3) and the review auto-promotes to \"Goals set\".",
            "You can edit or delete goals freely while the cycle is ACTIVE.",
        ],
        note="Write SMART goals: specific, measurable, attainable, relevant, time-bound.",
    ),
    Slide(
        title="Step 2 — Evaluate goals",
        subtitle="During the EVALUATION phase.",
        steps=[
            "Goals are now read-only. Each gets an outcome dropdown: Missed / Partial / Met / Exceeded.",
            "For quantitative goals, fill in the actual value.",
            "Add a comment per goal — what they did well, what they could do better.",
            "Save each as you go. The form remembers your work.",
        ],
    ),
    Slide(
        title="Step 3 — Overall rating + submit",
        subtitle="The summary that flows into bonus decisions.",
        steps=[
            "Below the goals: Submit review form.",
            "Pick the overall rating from the configured scale (e.g. 1–5: Below / Approaching / Meets / Exceeds / Outstanding).",
            "Write the manager narrative — strengths, growth areas, context.",
            "Fill any retail extras the cycle requires (sales target/actual, attendance, promotion-ready).",
            "Click Submit → review goes to the employee for acknowledgement, locked from edits.",
        ],
        note="Strict mode: all goals must be evaluated before you can submit.",
    ),
    Slide(
        title="Probation reviews",
        subtitle="A special 3-month flow for new hires.",
        bullets=[
            "No goals, no rating — just a confirmation decision.",
            "Decision options: CONFIRMED (pass probation) / EXTENDED (additional period) / NOT_CONFIRMED (terminate).",
            "Decision is required to submit. Legally meaningful in Malaysia — a confirmation letter must follow your decision.",
            "Notes field captures your rationale (kept on file).",
        ],
    ),
    # SECTION: Timesheet approvals
    Slide(title="Timesheet Approvals", subtitle="If you manage part-time staff.", section=True),
    Slide(
        title="The timesheet approval queue",
        screenshot="manager-time-approvals",
        bullets=[
            "Sidebar → Timesheet → Approvals tab.",
            "Submissions are grouped per employee, per week.",
            "Each week shows total hours + how many fall on public holidays (with auto-detection).",
            "\"Approve whole week\" approves all entries in one click, or do it per-entry.",
        ],
    ),
    Slide(
        title="Rejecting a timesheet entry",
        steps=[
            "On a specific entry, click \"Reject\".",
            "A comment box appears — enter a clear reason (this goes to the employee).",
            "Click \"Send rejection\" — entry returns to the employee as DRAFT, with your reason visible.",
            "They can fix and re-submit.",
        ],
        note="Use rejection sparingly. For repeated issues, talk to the employee directly first.",
    ),
    # SECTION: Other things
    Slide(title="Other day-to-day", subtitle="Calendars, documents, profile.", section=True),
    Slide(
        title="Team calendar",
        screenshot="manager-team-calendar",
        bullets=[
            "Shows who's out and when across your direct reports.",
            "Coloured by country (SG / MY).",
            "Useful when planning coverage or shift schedules.",
        ],
    ),
    Slide(
        title="Documents",
        screenshot="manager-documents",
        bullets=[
            "Read-only company documents (handbook, policies, holiday calendar).",
            "If HR uploaded employee-specific docs to one of your reports, you can read them.",
            "You can't upload — that's HR Admin's job.",
        ],
    ),
    Slide(
        title="Your own profile",
        bullets=[
            "Sidebar → Change Password — update your password any time.",
            "Your dashboard widgets are personal — they show YOUR pending approvals + birthdays.",
            "Need to update your phone / NRIC / bank info? Ask Admin for now.",
        ],
    ),
    Slide(
        title="That's the manager flow",
        subtitle="Walk through every step at your own pace.",
        bullets=[
            "Questions? Ask your HR Admin.",
            "The system enforces fair treatment by giving everyone the same workflow — but you bring the judgment.",
        ],
    ),
]


# ---- Employee deck ----
EMPLOYEE_SLIDES: List[Slide] = [
    Slide(
        title="Welcome — Employee guide",
        subtitle="How to use InsideHR for your own time off, expenses, and reviews.",
        bullets=[
            "InsideHR is the platform for your HR-related needs: applying for leave, claiming expenses, completing your performance review, and (for part-timers) logging your hours.",
            "Your access is limited to your own records — your manager sees your submissions; admin oversees the whole system.",
            "This deck walks each task step-by-step. ~20 slides; ~15 minutes.",
        ],
    ),
    Slide(
        title="Logging in for the first time",
        steps=[
            "Open the InsideHR URL provided by HR.",
            "Enter your work email + the temporary password from your welcome email.",
            "You'll be forced to change your password on first login (8+ characters, mix it up).",
            "After saving the new password, you land on your dashboard. That's home.",
        ],
        note="Lost your password? Click \"Forgot password\" on the login screen — you'll get a reset email.",
    ),
    Slide(
        title="Your dashboard",
        screenshot="employee-dashboard",
        bullets=[
            "Top-left widgets summarize your leave balances and any pending requests.",
            "Country holidays show the next few public holidays for your country.",
            "Birthday widget shows colleagues with upcoming birthdays — a nice reminder.",
            "Use the sidebar to navigate to specific features.",
        ],
    ),
    # SECTION: Leave
    Slide(title="Time Off (Leave)", subtitle="Applying, tracking, balances.", section=True),
    Slide(
        title="See your balances",
        screenshot="employee-leave",
        bullets=[
            "Annual leave, sick leave, maternity / paternity leave (if applicable), etc.",
            "Each card shows: available days, used days, pending days, and total entitlement.",
            "If you see a number you don't recognise, ask HR — balances can be manually adjusted.",
        ],
    ),
    Slide(
        title="Submitting a leave request",
        screenshot="employee-leave-request",
        steps=[
            "Sidebar → Time Off → \"Request leave\".",
            "Pick the leave type — only ones you have balance for show as available.",
            "Use the date picker. Working days are counted automatically (excludes weekends + your country's holidays).",
            "If it's a half-day, pick AM or PM.",
            "Add an optional reason + (if required) an attachment (e.g. MC for sick leave).",
            "Submit → request goes to your reporting manager for approval.",
        ],
        note="The system shows you a preview of working days deducted before you confirm.",
    ),
    Slide(
        title="Tracking your request",
        bullets=[
            "After submitting, your request shows on the Time Off page with a status pill.",
            "Pending → Approved (deducted from balance) or Rejected (with manager's reason).",
            "You can cancel a still-pending request at any time.",
            "Approved leave appears on the team calendar so colleagues know.",
        ],
    ),
    # SECTION: Expenses
    Slide(title="Expenses", subtitle="Claiming reimbursable spend.", section=True),
    Slide(
        title="Submitting an expense",
        screenshot="employee-expense-new",
        steps=[
            "Sidebar → Expenses → \"+ New expense\".",
            "Pick a category (Local Transport / Meals / Office / etc.).",
            "Amount + currency. Receipt date. Merchant name.",
            "Add a description if helpful.",
            "Upload the receipt — PDF, JPG, or PNG.",
            "Save as Draft, then click \"Submit for approval\" when ready.",
        ],
        note="Receipts live securely in Google Drive; only you, your manager, and admin can see them.",
    ),
    Slide(
        title="Tracking your expenses",
        screenshot="employee-expenses",
        bullets=[
            "Tabs at top: Draft / For Approval / Approved / Rejected / Reimbursed.",
            "Each row shows the status pill at a glance.",
            "Click into one to see history: when submitted, who approved, when reimbursed.",
            "Rejected expenses come back editable, with the manager's reason.",
        ],
    ),
    # SECTION: Performance Review
    Slide(title="Your Performance Reviews", subtitle="What to expect from a review cycle.", section=True),
    Slide(
        title="How review cycles work for you",
        bullets=[
            "Admin runs review cycles every quarter / half-year / year, depending on company practice.",
            "Your reporting manager writes your goals at the start, then evaluates you at the end.",
            "You don't fill out goals yourself in v1 — your manager does. You'll see them as read-only.",
            "At the end, you read your manager's evaluation and \"Acknowledge\" — that locks the review on record.",
        ],
        note="If a probation review applies (within your first 3 months), you'll see CONFIRMED / EXTENDED / NOT_CONFIRMED there instead of a rating.",
    ),
    Slide(
        title="Seeing your reviews",
        screenshot="employee-perf-me",
        bullets=[
            "Sidebar → Performance → My reviews.",
            "Each row = one review cycle. Status tells you what's happening now.",
            "Click \"View\" to see the goals + your manager's evaluation when ready.",
        ],
    ),
    Slide(
        title="Acknowledging your review",
        screenshot="employee-perf-review-detail",
        steps=[
            "When your review is \"Awaiting acknowledgement\", click \"View →\" or \"Acknowledge →\".",
            "Read each goal + outcome + manager comment.",
            "Read your overall rating + the manager narrative.",
            "Optionally add a comment (agreement, disagreement, context — kept on record).",
            "Click \"Acknowledge review\" → it's now locked. Done.",
        ],
        note="Acknowledgement isn't agreement — it's just confirmation you've seen it. Disagree? Add a comment and talk to your manager.",
    ),
    # SECTION: Timesheet (only relevant for PT staff)
    Slide(title="Timesheet (Part-time staff only)", subtitle="Logging your hours daily and submitting weekly.", section=True),
    Slide(
        title="Logging your daily hours",
        screenshot="pt-timesheet",
        steps=[
            "Sidebar → Timesheet → you'll see the current week as 7 day-cards.",
            "Click any day → enter hours worked, optional start/end time, optional break minutes, optional note.",
            "Save → that day's card shows \"Draft\" + the hours you logged.",
            "Repeat for each day you worked that week.",
        ],
        note="Public holidays auto-detect from your country's calendar; if you worked on one, you'll get 2× pay (3× for overtime).",
    ),
    Slide(
        title="Submitting your week",
        steps=[
            "Once all your worked days are logged, click \"Submit week\" at the top right.",
            "All your drafts lock as \"Submitted\" — your manager gets them in their approval queue.",
            "Wait for approval. Approved entries flow into the monthly payroll.",
            "If your manager rejects an entry, it comes back as \"Draft\" with their reason — fix and resubmit.",
        ],
        note="You can edit drafts and rejected entries up to 14 days back. Submitted entries are read-only until approved or rejected.",
    ),
    # Misc
    Slide(title="Other useful things", section=True, subtitle="Calendar, docs, and your own profile."),
    Slide(
        title="Team calendar",
        screenshot="employee-team-calendar",
        bullets=[
            "See who's out across your department / company.",
            "Useful for planning meetings or coverage.",
            "Coloured by country.",
        ],
    ),
    Slide(
        title="Documents",
        screenshot="employee-documents",
        bullets=[
            "Two sections: Company documents (handbook, policies, holiday calendar) + My documents (your offer letter, contract, etc.).",
            "Click to download — files are stored in Google Drive.",
            "Can't see something? Ask HR — admin controls who sees what.",
        ],
    ),
    Slide(
        title="Your password",
        bullets=[
            "Sidebar bottom → Change Password — update any time.",
            "Forgot your password? Click \"Forgot password\" on the login screen and you'll get a reset email.",
            "Never share your password with anyone, even HR. They can reset it but not see it.",
        ],
    ),
    Slide(
        title="That's everything",
        subtitle="Welcome to InsideHR.",
        bullets=[
            "Questions about your records? Ask your manager first, then HR.",
            "Found a bug or have a feature request? Let HR know — they can pass it up.",
        ],
    ),
]


# ============================================================
# Deck assembly
# ============================================================


def build_deck(role: str, sub: str, slides: List[Slide], filename: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs, role=role, sub=sub)

    page = 2
    for spec in slides:
        if spec.section:
            build_section_slide(prs, title=spec.title, subtitle=spec.subtitle or "")
        else:
            build_content_slide(prs, spec, deck_name=f"InsideHR · {role} Training", page_num=page)
        page += 1

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out_path = OUT_DIR / filename
    prs.save(out_path)
    print(f"  ✓ {filename} ({len(slides) + 1} slides) → {out_path}")


def main():
    print("Generating InsideHR training decks…")
    build_deck("Admin", "The full picture, end to end.", ADMIN_SLIDES, "admin.pptx")
    build_deck("Manager", "Goal-setting, evaluation, approvals.", MANAGER_SLIDES, "manager.pptx")
    build_deck("Employee", "Self-service: leave, expenses, reviews, timesheet.", EMPLOYEE_SLIDES, "employee.pptx")
    print(f"\nDone. Decks in {OUT_DIR}")


if __name__ == "__main__":
    main()
